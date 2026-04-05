#!/usr/bin/env python3
"""Generate VModelWorkflow stakeholder presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)       # Deep navy
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xD6)    # Bright blue
ACCENT_TEAL = RGBColor(0x00, 0xB4, 0xA0)    # Teal
LIGHT_TEXT = RGBColor(0xFF, 0xFF, 0xFF)      # White
SUBTLE_TEXT = RGBColor(0xB0, 0xB8, 0xC8)     # Light gray-blue
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)       # Dark navy
WARM_BG = RGBColor(0xF5, 0xF7, 0xFA)        # Off-white
ORANGE = RGBColor(0xFF, 0x8C, 0x42)          # Warm orange accent
RED_ACCENT = RGBColor(0xE8, 0x4D, 0x4D)     # Red for pain points
GREEN_ACCENT = RGBColor(0x2E, 0xCC, 0x71)    # Green for solutions

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_bg(slide, color):
    """Set solid background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, alpha=None):
    """Add a colored rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=LIGHT_TEXT, bold=False, align=PP_ALIGN.LEFT,
                 font_name="Calibri", line_spacing=1.2):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_multiline_box(slide, left, top, width, height, lines, font_size=16,
                      color=LIGHT_TEXT, bold=False, align=PP_ALIGN.LEFT,
                      font_name="Calibri", line_spacing=1.3, bullet=False):
    """Add text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if bullet and line.strip():
            p.text = line
        else:
            p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = align
        p.space_after = Pt(4)
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_accent_bar(slide, left, top, width=Inches(0.08), height=Inches(0.8), color=ACCENT_BLUE):
    """Add a thin accent bar."""
    return add_shape(slide, left, top, width, height, color)


def add_bottom_bar(slide):
    """Add bottom accent stripe."""
    add_shape(slide, Inches(0), SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), ACCENT_BLUE)


def add_slide_number(slide, num, total):
    add_text_box(slide, SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.5),
                 Inches(1), Inches(0.4), f"{num}/{total}",
                 font_size=11, color=SUBTLE_TEXT, align=PP_ALIGN.RIGHT)


TOTAL_SLIDES = 13

# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide, DARK_BG)

# Decorative top gradient bar
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.04), ACCENT_BLUE)
add_shape(slide, Inches(0), Inches(0.04), SLIDE_W, Inches(0.02), ACCENT_TEAL)

# Title
add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.2),
             "VModelWorkflow", font_size=54, color=LIGHT_TEXT, bold=True,
             font_name="Calibri Light", align=PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(0.8),
             "V-Model Development Framework for Safety-Critical AI", font_size=24,
             color=ACCENT_TEAL, align=PP_ALIGN.CENTER, font_name="Calibri")

# Divider line
add_shape(slide, Inches(5.5), Inches(4.2), Inches(2.3), Inches(0.03), ACCENT_BLUE)

# Tagline
add_text_box(slide, Inches(2), Inches(4.6), Inches(9.3), Inches(1.2),
             "Bridging safety standards, clean code, and AI-driven development",
             font_size=18, color=SUBTLE_TEXT, align=PP_ALIGN.CENTER)

add_bottom_bar(slide)
add_slide_number(slide, 1, TOTAL_SLIDES)

# ============================================================
# SLIDE 2: The Problem
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
             "The Problem", font_size=36, color=LIGHT_TEXT, bold=True)
add_shape(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.04), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.8),
             "Safety-critical software must follow strict V-model standards. Compliance is expensive, painful, and often disconnected from actual quality.",
             font_size=17, color=SUBTLE_TEXT)

# Pain point cards
pain_points = [
    ("Understanding Standards", "Dense, domain-specific documents.\nTranslating to daily work is hard."),
    ("No Practical Frameworks", "Tools handle pieces, but none\nprovide end-to-end coverage."),
    ("Documentation Overhead", "Traceability eats 30-50% of\ndocumentation time."),
    ("Quality Disconnect", "Compliance becomes a checkbox\nexercise, not engineering quality."),
]

card_w = Inches(2.7)
card_h = Inches(2.8)
start_x = Inches(0.8)
gap = Inches(0.3)

for i, (title, desc) in enumerate(pain_points):
    x = start_x + i * (card_w + gap)
    y = Inches(2.8)
    card = add_rounded_rect(slide, x, y, card_w, card_h, RGBColor(0x24, 0x24, 0x40))
    # Red top accent
    add_shape(slide, x + Inches(0.3), y + Inches(0.3), Inches(0.5), Inches(0.05), RED_ACCENT)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.5), card_w - Inches(0.6), Inches(0.6),
                 title, font_size=16, color=LIGHT_TEXT, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.2), card_w - Inches(0.6), Inches(1.4),
                 desc, font_size=14, color=SUBTLE_TEXT)

# Result callout
add_text_box(slide, Inches(0.8), Inches(6.0), Inches(11), Inches(0.6),
             "Result: significant overhead, audit failures, and compliance work that doesn't improve actual software.",
             font_size=15, color=ORANGE, bold=True)

add_bottom_bar(slide)
add_slide_number(slide, 2, TOTAL_SLIDES)

# ============================================================
# SLIDE 3: What is the V-Model?
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
             "What is the V-Model?", font_size=36, color=LIGHT_TEXT, bold=True)
add_shape(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.04), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
             "A development & verification framework used across all safety-critical industries.",
             font_size=17, color=SUBTLE_TEXT)

# V-diagram using text-based representation
v_levels = [
    ("System Requirements", "System Tests"),
    ("SW Requirements", "Integration Tests"),
    ("SW Architecture", "Component Tests"),
    ("Detailed Design", "Unit Tests"),
]

center_x = SLIDE_W / 2
start_y = Inches(2.4)
level_h = Inches(0.85)

for i, (left_label, right_label) in enumerate(v_levels):
    y = start_y + i * level_h
    offset = Inches(0.7) * i

    # Left side box
    lx = Inches(1.2) + offset
    box_w = Inches(2.6)
    box = add_rounded_rect(slide, lx, y, box_w, Inches(0.6), RGBColor(0x00, 0x70, 0xA0))
    add_text_box(slide, lx, y + Inches(0.1), box_w, Inches(0.45),
                 left_label, font_size=14, color=LIGHT_TEXT, bold=True, align=PP_ALIGN.CENTER)

    # Right side box
    rx = SLIDE_W - Inches(1.2) - offset - box_w
    box = add_rounded_rect(slide, rx, y, box_w, Inches(0.6), RGBColor(0x00, 0x8A, 0x7A))
    add_text_box(slide, rx, y + Inches(0.1), box_w, Inches(0.45),
                 right_label, font_size=14, color=LIGHT_TEXT, bold=True, align=PP_ALIGN.CENTER)

    # Traceability arrow label
    add_text_box(slide, lx + box_w + Inches(0.15), y + Inches(0.1),
                 rx - lx - box_w - Inches(0.3), Inches(0.4),
                 "< - - traceability - - >", font_size=11, color=ACCENT_TEAL,
                 align=PP_ALIGN.CENTER)

# Source Code at bottom center
sc_w = Inches(2.2)
sc_x = center_x - sc_w / 2
sc_y = start_y + 4 * level_h
box = add_rounded_rect(slide, sc_x, sc_y, sc_w, Inches(0.6), ORANGE)
add_text_box(slide, sc_x, sc_y + Inches(0.1), sc_w, Inches(0.45),
             "Source Code", font_size=14, color=DARK_TEXT, bold=True, align=PP_ALIGN.CENTER)

# Footer note
add_text_box(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.5),
             "Universal across DO-178C (aviation) | ISO 26262 (automotive) | ASPICE | IEC 62304 (medical) | EN 50128 (railway)",
             font_size=13, color=SUBTLE_TEXT, align=PP_ALIGN.CENTER)

add_bottom_bar(slide)
add_slide_number(slide, 3, TOTAL_SLIDES)

# ============================================================
# SLIDE 4: The Opportunity — AI Changes Everything
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
             "The Opportunity: AI Changes Everything", font_size=36, color=LIGHT_TEXT, bold=True)
add_shape(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.04), ACCENT_TEAL)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
             "AI performs better with more guardrails, stricter instructions, and layered design. This is exactly what the V-model provides.",
             font_size=18, color=ACCENT_TEAL, bold=True)

# Three-column convergence table
headers = ["What Standards Demand", "What Clean Code Teaches", "What AI Needs"]
header_colors = [ACCENT_BLUE, ACCENT_TEAL, ORANGE]
rows = [
    ["Design before code", "Understand before coding", "Spec-first approach"],
    ["Small, testable units", "Functions under 20 lines", "Accuracy degrades with size"],
    ["Naming conventions", "Names reveal intent", "Names are primary context"],
    ["No dead code", "Delete unused code", "AI over-generates; prune"],
    ["Structured error handling", "No null returns", "AI omits error paths"],
    ["Testability", "If hard to test, redesign", "Test-first = success criteria"],
]

col_w = Inches(3.6)
table_start_x = Inches(0.8)
table_start_y = Inches(2.5)
row_h = Inches(0.55)

for j, (header, hcolor) in enumerate(zip(headers, header_colors)):
    x = table_start_x + j * (col_w + Inches(0.2))
    # Header
    hbox = add_rounded_rect(slide, x, table_start_y, col_w, Inches(0.55), hcolor)
    add_text_box(slide, x, table_start_y + Inches(0.1), col_w, Inches(0.4),
                 header, font_size=14, color=LIGHT_TEXT, bold=True, align=PP_ALIGN.CENTER)
    # Rows
    for i, row in enumerate(rows):
        ry = table_start_y + Inches(0.65) + i * row_h
        bg_color = RGBColor(0x22, 0x22, 0x3A) if i % 2 == 0 else RGBColor(0x28, 0x28, 0x44)
        add_shape(slide, x, ry, col_w, row_h - Inches(0.03), bg_color)
        add_text_box(slide, x + Inches(0.15), ry + Inches(0.1), col_w - Inches(0.3), row_h,
                     row[j], font_size=13, color=SUBTLE_TEXT, align=PP_ALIGN.CENTER)

# Conclusion
add_text_box(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.6),
             "V-model compliance, clean code, and AI-driven development demand the same discipline -- from different motivations.",
             font_size=16, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

add_bottom_bar(slide)
add_slide_number(slide, 4, TOTAL_SLIDES)

# ============================================================
# SLIDE 5: What VModelWorkflow Provides (Overview)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
             "Four Independent Components", font_size=36, color=LIGHT_TEXT, bold=True)
add_shape(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.04), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
             "Each component works alone. The real value is the combination.",
             font_size=17, color=SUBTLE_TEXT)

components = [
    ("1", "Documentation", "Foundation", "V-model education, best practices,\nanti-patterns, examples.\nDepth equivalent to major standards.", ACCENT_BLUE),
    ("2", "Templates & Schemas", "Structure", "Artifact definitions, checklists,\nassurance levels, domain\ntranslation across standards.", ACCENT_TEAL),
    ("3", "Traceability", "Validation", "Many-to-many linking, automated\nvalidation, orphan detection,\nchange impact analysis.", ORANGE),
    ("4", "AI Skills", "Automation", "Craft skills (standalone) +\nFramework skills (orchestration).\nWorks on smaller AI models.", RGBColor(0xAA, 0x66, 0xCC)),
]

card_w = Inches(2.7)
card_h = Inches(3.5)
start_x = Inches(0.8)
gap = Inches(0.3)

for i, (num, title, subtitle, desc, color) in enumerate(components):
    x = start_x + i * (card_w + gap)
    y = Inches(2.3)

    card = add_rounded_rect(slide, x, y, card_w, card_h, RGBColor(0x24, 0x24, 0x40))

    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), y + Inches(0.3),
                                     Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_text_box(slide, x + Inches(0.3), y + Inches(0.35), Inches(0.55), Inches(0.45),
                 num, font_size=20, color=LIGHT_TEXT, bold=True, align=PP_ALIGN.CENTER)

    add_text_box(slide, x + Inches(1.0), y + Inches(0.35), card_w - Inches(1.3), Inches(0.4),
                 title, font_size=16, color=LIGHT_TEXT, bold=True)
    add_text_box(slide, x + Inches(1.0), y + Inches(0.7), card_w - Inches(1.3), Inches(0.3),
                 subtitle, font_size=12, color=color)

    add_shape(slide, x + Inches(0.3), y + Inches(1.15), card_w - Inches(0.6), Inches(0.02), color)

    add_text_box(slide, x + Inches(0.3), y + Inches(1.4), card_w - Inches(0.6), Inches(2.0),
                 desc, font_size=13, color=SUBTLE_TEXT)

add_bottom_bar(slide)
add_slide_number(slide, 5, TOTAL_SLIDES)

# ============================================================
# SLIDE 6: How It Works Together
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
             "How It Works Together", font_size=36, color=LIGHT_TEXT, bold=True)
add_shape(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.04), ACCENT_BLUE)

# Adoption table
adopt_rows = [
    ("Documentation alone", "Comprehensive V-model education & best practices"),
    ("Templates alone", "Artifact definitions & checklists for manual compliance"),
    ("Craft skills alone", "Quality writing, test derivation, code dev for any project"),
    ("Docs + Templates + Traceability", "Manual V-model dev with automated completeness checking"),
    ("All four together", "Full agentic V-model development"),
]

table_y = Inches(1.8)
row_h = Inches(0.75)
col1_w = Inches(4.5)
col2_w = Inches(7.5)

# Header
add_shape(slide, Inches(0.8), table_y, col1_w, Inches(0.6), ACCENT_BLUE)
add_text_box(slide, Inches(0.8), table_y + Inches(0.12), col1_w, Inches(0.45),
             "What You Use", font_size=15, color=LIGHT_TEXT, bold=True, align=PP_ALIGN.CENTER)
add_shape(slide, Inches(0.8) + col1_w + Inches(0.05), table_y, col2_w, Inches(0.6), ACCENT_BLUE)
add_text_box(slide, Inches(0.8) + col1_w + Inches(0.05), table_y + Inches(0.12), col2_w, Inches(0.45),
             "What You Get", font_size=15, color=LIGHT_TEXT, bold=True, align=PP_ALIGN.CENTER)

for i, (use, get) in enumerate(adopt_rows):
    ry = table_y + Inches(0.7) + i * row_h
    bg = RGBColor(0x22, 0x22, 0x3A) if i % 2 == 0 else RGBColor(0x28, 0x28, 0x44)
    is_last = (i == len(adopt_rows) - 1)

    add_shape(slide, Inches(0.8), ry, col1_w, row_h - Inches(0.05), bg)
    tc = ORANGE if is_last else SUBTLE_TEXT
    add_text_box(slide, Inches(1.0), ry + Inches(0.15), col1_w - Inches(0.4), row_h,
                 use, font_size=14, color=tc, bold=is_last)

    add_shape(slide, Inches(0.8) + col1_w + Inches(0.05), ry, col2_w, row_h - Inches(0.05), bg)
    gc = ORANGE if is_last else LIGHT_TEXT
    add_text_box(slide, Inches(0.8) + col1_w + Inches(0.25), ry + Inches(0.15), col2_w - Inches(0.4), row_h,
                 get, font_size=14, color=gc, bold=is_last)

add_bottom_bar(slide)
add_slide_number(slide, 6, TOTAL_SLIDES)

# ============================================================
# SLIDE 7: The Workflow
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
             "The Workflow: Human + AI Collaboration", font_size=36, color=LIGHT_TEXT, bold=True)
add_shape(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.04), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
             "For each V-model layer, a mid-senior engineer orchestrates AI through three phases:",
             font_size=17, color=SUBTLE_TEXT)

# Three phase boxes
phases = [
    ("Research & Plan", "HUMAN-DRIVEN", [
        "Human provides context",
        "AI gathers & analyzes",
        "Back-and-forth discussion",
        "Agreed plan produced",
    ], ACCENT_BLUE),
    ("Implementation", "AGENT-ORCHESTRATED", [
        "AI writes artifacts",
        "AI self-checks quality",
        "Review agent validates",
        "Traceability auto-updated",
    ], ACCENT_TEAL),
    ("Final Review", "HUMAN-DRIVEN", [
        "Human reviews output",
        "Approves or rejects",
        "Quality gate enforced",
        "Transitions to next layer",
    ], ORANGE),
]

phase_w = Inches(3.5)
phase_h = Inches(3.8)
phase_gap = Inches(0.4)
phase_start_x = Inches(0.8)

for i, (title, subtitle, items, color) in enumerate(phases):
    x = phase_start_x + i * (phase_w + phase_gap)
    y = Inches(2.3)

    card = add_rounded_rect(slide, x, y, phase_w, phase_h, RGBColor(0x24, 0x24, 0x40))

    # Phase header bar
    add_shape(slide, x, y, phase_w, Inches(0.06), color)

    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), phase_w - Inches(0.6), Inches(0.4),
                 subtitle, font_size=11, color=color, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.7), phase_w - Inches(0.6), Inches(0.5),
                 title, font_size=20, color=LIGHT_TEXT, bold=True)

    add_shape(slide, x + Inches(0.3), y + Inches(1.2), phase_w - Inches(0.6), Inches(0.02), color)

    for j, item in enumerate(items):
        iy = y + Inches(1.5) + j * Inches(0.5)
        # Bullet dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.4), iy + Inches(0.08),
                                      Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        add_text_box(slide, x + Inches(0.7), iy, phase_w - Inches(1.0), Inches(0.4),
                     item, font_size=14, color=SUBTLE_TEXT)

    # Arrow between phases
    if i < 2:
        ax = x + phase_w + Inches(0.05)
        ay = y + phase_h / 2 - Inches(0.15)
        add_text_box(slide, ax, ay, Inches(0.3), Inches(0.3),
                     ">", font_size=24, color=SUBTLE_TEXT, align=PP_ALIGN.CENTER, bold=True)

add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.5),
             "Not an autonomous pipeline. Human-orchestrated, AI-augmented.",
             font_size=15, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

add_bottom_bar(slide)
add_slide_number(slide, 7, TOTAL_SLIDES)

# ============================================================
# SLIDE 8: DRTDD Innovation
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
             "The Innovation: DRTDD", font_size=36, color=LIGHT_TEXT, bold=True)
add_shape(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.04), ACCENT_TEAL)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
             "Design-Requirement-Test Driven Development extends TDD with V-model discipline.",
             font_size=17, color=SUBTLE_TEXT)

# DRTDD flow
steps = [
    ("REQUIRE", ACCENT_BLUE),
    ("DESIGN", ACCENT_BLUE),
    ("TEST\n(red)", RED_ACCENT),
    ("IMPLEMENT\n(green)", GREEN_ACCENT),
    ("REFACTOR", ACCENT_TEAL),
    ("VERIFY", ORANGE),
]

step_w = Inches(1.7)
step_h = Inches(1.1)
flow_start_x = Inches(0.8)
flow_y = Inches(2.6)
arrow_w = Inches(0.25)

for i, (label, color) in enumerate(steps):
    x = flow_start_x + i * (step_w + arrow_w)
    box = add_rounded_rect(slide, x, flow_y, step_w, step_h, color)
    add_text_box(slide, x, flow_y + Inches(0.2), step_w, step_h - Inches(0.2),
                 label, font_size=16, color=LIGHT_TEXT, bold=True, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        add_text_box(slide, x + step_w, flow_y + Inches(0.25), arrow_w, step_h,
                     ">", font_size=22, color=SUBTLE_TEXT, align=PP_ALIGN.CENTER, bold=True)

# Feedback loop label
add_text_box(slide, Inches(0.8), Inches(4.0), Inches(11.5), Inches(0.5),
             "< - - - - - - - - - - - -  gap found? feed back to REQUIRE  - - - - - - - - - - - - >",
             font_size=13, color=RED_ACCENT, align=PP_ALIGN.CENTER)

# Benefits below
benefits_title_y = Inches(4.8)
add_text_box(slide, Inches(0.8), benefits_title_y, Inches(11), Inches(0.5),
             "What DRTDD adds beyond standard TDD:", font_size=16, color=LIGHT_TEXT, bold=True)

benefits = [
    "Requirements traceability at every step",
    "Design rationale captured as artifacts",
    "Safety analysis integration",
    "Verification completeness beyond code coverage",
    "Human gates between phases",
    "Feedback loops to requirement level, not just code",
]

for i, b in enumerate(benefits):
    col = i // 3
    row = i % 3
    bx = Inches(1.2) + col * Inches(5.5)
    by = benefits_title_y + Inches(0.6) + row * Inches(0.45)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, bx, by + Inches(0.08),
                                  Inches(0.1), Inches(0.1))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT_TEAL
    dot.line.fill.background()
    add_text_box(slide, bx + Inches(0.25), by, Inches(5), Inches(0.4),
                 b, font_size=14, color=SUBTLE_TEXT)

add_bottom_bar(slide)
add_slide_number(slide, 8, TOTAL_SLIDES)

# ============================================================
# SLIDE 9: Market Entry — Legacy Retrofit
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
             "Primary Market Entry: Legacy Retrofit", font_size=36, color=LIGHT_TEXT, bold=True)
add_shape(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.04), ORANGE)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.8),
             "The biggest value is not greenfield -- it's making existing legacy codebases V-model compliant.\nThe volume of legacy code needing compliance far exceeds new development.",
             font_size=17, color=SUBTLE_TEXT)

capabilities = [
    ("Analyze Existing Code", "Understand structure, dependencies,\nimplicit design decisions"),
    ("Reverse-Engineer Artifacts", "Infer requirements, designs, and\ntest coverage from code"),
    ("Identify Gaps", "What's missing, implicit, or needs\nformalization"),
    ("Suggest Improvements", "Highlight where code falls short\nof best practices"),
]

cap_w = Inches(5.3)
cap_h = Inches(1.1)
cap_start_y = Inches(2.8)

for i, (title, desc) in enumerate(capabilities):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + col * (cap_w + Inches(0.4))
    y = cap_start_y + row * (cap_h + Inches(0.3))

    card = add_rounded_rect(slide, x, y, cap_w, cap_h, RGBColor(0x24, 0x24, 0x40))
    add_accent_bar(slide, x + Inches(0.15), y + Inches(0.2), Inches(0.06), Inches(0.7), ORANGE)
    add_text_box(slide, x + Inches(0.5), y + Inches(0.15), cap_w - Inches(0.7), Inches(0.35),
                 title, font_size=15, color=LIGHT_TEXT, bold=True)
    add_text_box(slide, x + Inches(0.5), y + Inches(0.5), cap_w - Inches(0.7), Inches(0.55),
                 desc, font_size=13, color=SUBTLE_TEXT)

add_text_box(slide, Inches(0.8), Inches(5.6), Inches(11.5), Inches(0.5),
             "Same components (documentation, templates, traceability, skills) applied bottom-up instead of top-down.",
             font_size=15, color=ACCENT_TEAL, bold=True, align=PP_ALIGN.CENTER)

add_bottom_bar(slide)
add_slide_number(slide, 9, TOTAL_SLIDES)

# ============================================================
# SLIDE 10: Market Context & Standards
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
             "Market Context", font_size=36, color=LIGHT_TEXT, bold=True)
add_shape(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.04), ACCENT_BLUE)

# Standards table
std_headers = ["Domain", "Standard", "Safety Levels"]
std_rows = [
    ["Aviation", "DO-178C / DO-330", "DAL A-E"],
    ["Automotive", "ISO 26262", "ASIL A-D"],
    ["Automotive", "ASPICE", "Level 1-5"],
    ["Medical", "IEC 62304", "Class A-C"],
    ["Railway", "EN 50128", "SIL 0-4"],
]

table_x = Inches(0.8)
table_y = Inches(1.6)
col_widths = [Inches(2), Inches(3), Inches(2)]
row_h = Inches(0.5)

# Header
hx = table_x
for j, (header, cw) in enumerate(zip(std_headers, col_widths)):
    add_shape(slide, hx, table_y, cw - Inches(0.05), row_h, ACCENT_BLUE)
    add_text_box(slide, hx + Inches(0.15), table_y + Inches(0.08), cw - Inches(0.3), row_h,
                 header, font_size=13, color=LIGHT_TEXT, bold=True, align=PP_ALIGN.CENTER)
    hx += cw

for i, row in enumerate(std_rows):
    ry = table_y + Inches(0.55) + i * row_h
    bg = RGBColor(0x22, 0x22, 0x3A) if i % 2 == 0 else RGBColor(0x28, 0x28, 0x44)
    rx = table_x
    for j, (cell, cw) in enumerate(zip(row, col_widths)):
        add_shape(slide, rx, ry, cw - Inches(0.05), row_h - Inches(0.03), bg)
        add_text_box(slide, rx + Inches(0.15), ry + Inches(0.08), cw - Inches(0.3), row_h,
                     cell, font_size=12, color=SUBTLE_TEXT, align=PP_ALIGN.CENTER)
        rx += cw

# AI in Regulated Industries - right side
ai_x = Inches(7.5)
add_text_box(slide, ai_x, Inches(1.6), Inches(5.5), Inches(0.4),
             "AI in Safety-Critical (2026)", font_size=16, color=LIGHT_TEXT, bold=True)
add_shape(slide, ai_x, Inches(2.0), Inches(1.5), Inches(0.03), ACCENT_TEAL)

ai_facts = [
    "No LLM tool qualified under DO-330 yet",
    "~60% of DAL-A/B orgs ban generative AI",
    "~35% of automotive OEMs cautiously piloting",
    "Documentation gen is most accepted use case",
]

for i, fact in enumerate(ai_facts):
    fy = Inches(2.3) + i * Inches(0.5)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, ai_x + Inches(0.1), fy + Inches(0.08),
                                  Inches(0.1), Inches(0.1))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT_TEAL
    dot.line.fill.background()
    add_text_box(slide, ai_x + Inches(0.35), fy, Inches(5), Inches(0.4),
                 fact, font_size=13, color=SUBTLE_TEXT)

# Gap statement
gap_y = Inches(4.7)
gap_box = add_rounded_rect(slide, Inches(0.8), gap_y, Inches(11.5), Inches(1.2),
                            RGBColor(0x24, 0x24, 0x40))
add_shape(slide, Inches(0.8), gap_y, Inches(0.08), Inches(1.2), ORANGE)
add_text_box(slide, Inches(1.2), gap_y + Inches(0.15), Inches(10.8), Inches(0.4),
             "The Gap VModelWorkflow Fills", font_size=16, color=ORANGE, bold=True)
add_text_box(slide, Inches(1.2), gap_y + Inches(0.55), Inches(10.8), Inches(0.6),
             "No existing framework addresses agentic AI in safety-critical development. VModelWorkflow treats AI as a suggestion engine -- AI generates, deterministic tools verify, humans approve.",
             font_size=14, color=SUBTLE_TEXT)

add_bottom_bar(slide)
add_slide_number(slide, 10, TOTAL_SLIDES)

# ============================================================
# SLIDE 11: Competitive Landscape
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
             "Competitive Landscape", font_size=36, color=LIGHT_TEXT, bold=True)
add_shape(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.04), ACCENT_BLUE)

comp_headers = ["Tool", "What It Does", "What It Misses"]
comp_rows = [
    ["DOORS / Polarion", "Requirements management", "No AI, no code-level traceability"],
    ["Doorstop / StrictDoc", "Requirements-as-code", "No test derivation, limited validation"],
    ["TRLC + LOBSTER\n(BMW)", "Traceability with\nISO 26262 support", "No AI skills, no templates,\nnot end-to-end"],
    ["SonarQube /\nSpotBugs", "Code quality analysis", "No V-model awareness,\nno requirements link"],
    ["Copilot / Cursor", "AI code generation", "No safety awareness,\nno traceability"],
]

comp_col_widths = [Inches(2.5), Inches(3.5), Inches(4.5)]
comp_y = Inches(1.6)
comp_row_h = Inches(0.85)

# Headers
cx = Inches(0.8)
for j, (h, cw) in enumerate(zip(comp_headers, comp_col_widths)):
    add_shape(slide, cx, comp_y, cw - Inches(0.05), Inches(0.55), ACCENT_BLUE)
    add_text_box(slide, cx + Inches(0.15), comp_y + Inches(0.1), cw - Inches(0.3), Inches(0.4),
                 h, font_size=14, color=LIGHT_TEXT, bold=True, align=PP_ALIGN.CENTER)
    cx += cw

for i, row in enumerate(comp_rows):
    ry = comp_y + Inches(0.6) + i * comp_row_h
    bg = RGBColor(0x22, 0x22, 0x3A) if i % 2 == 0 else RGBColor(0x28, 0x28, 0x44)
    cx = Inches(0.8)
    for j, (cell, cw) in enumerate(zip(row, comp_col_widths)):
        add_shape(slide, cx, ry, cw - Inches(0.05), comp_row_h - Inches(0.05), bg)
        tc = RED_ACCENT if j == 2 else SUBTLE_TEXT
        add_text_box(slide, cx + Inches(0.15), ry + Inches(0.1), cw - Inches(0.3), comp_row_h,
                     cell, font_size=12, color=tc, align=PP_ALIGN.CENTER)
        cx += cw

# Bottom callout
add_rounded_rect(slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.7), RGBColor(0x24, 0x24, 0x40))
add_text_box(slide, Inches(1.2), Inches(6.1), Inches(10.8), Inches(0.5),
             "No competitor provides end-to-end DRTDD + V-model compliance + traceability + agentic orchestration.",
             font_size=15, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

add_bottom_bar(slide)
add_slide_number(slide, 11, TOTAL_SLIDES)

# ============================================================
# SLIDE 12: Estimated Impact
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
             "Estimated Impact", font_size=36, color=LIGHT_TEXT, bold=True)
add_shape(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.04), GREEN_ACCENT)

impacts = [
    ("Traceability\nMaintenance", "30-50% of\ndoc time", "Automated\nnear zero", "~90%"),
    ("Review\nOverhead", "+20-40%\ndev time", "AI pre-review\nreduces burden", "~60%"),
    ("Legacy\nRetrofit", "2-3x greenfield\neffort", "Incremental\nAI analysis", "~50%"),
    ("Audit\nPreparation", "Weeks of\ngathering", "Continuous\nevidence", "~80%"),
]

imp_w = Inches(2.7)
imp_h = Inches(3.8)
imp_start_x = Inches(0.8)
imp_gap = Inches(0.3)

for i, (area, current, improved, reduction) in enumerate(impacts):
    x = imp_start_x + i * (imp_w + imp_gap)
    y = Inches(1.8)

    card = add_rounded_rect(slide, x, y, imp_w, imp_h, RGBColor(0x24, 0x24, 0x40))

    add_text_box(slide, x + Inches(0.2), y + Inches(0.2), imp_w - Inches(0.4), Inches(0.8),
                 area, font_size=16, color=LIGHT_TEXT, bold=True, align=PP_ALIGN.CENTER)

    add_shape(slide, x + Inches(0.3), y + Inches(1.1), imp_w - Inches(0.6), Inches(0.02), SUBTLE_TEXT)

    add_text_box(slide, x + Inches(0.2), y + Inches(1.2), imp_w - Inches(0.4), Inches(0.3),
                 "Current", font_size=11, color=RED_ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), y + Inches(1.5), imp_w - Inches(0.4), Inches(0.6),
                 current, font_size=13, color=SUBTLE_TEXT, align=PP_ALIGN.CENTER)

    add_text_box(slide, x + Inches(0.2), y + Inches(2.2), imp_w - Inches(0.4), Inches(0.3),
                 "With VModelWorkflow", font_size=11, color=GREEN_ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), y + Inches(2.5), imp_w - Inches(0.4), Inches(0.6),
                 improved, font_size=13, color=SUBTLE_TEXT, align=PP_ALIGN.CENTER)

    # Big reduction number
    add_text_box(slide, x + Inches(0.2), y + Inches(3.1), imp_w - Inches(0.4), Inches(0.6),
                 reduction, font_size=28, color=GREEN_ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), y + Inches(3.55), imp_w - Inches(0.4), Inches(0.3),
                 "reduction", font_size=10, color=SUBTLE_TEXT, align=PP_ALIGN.CENTER)

# Overall
add_rounded_rect(slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.7), RGBColor(0x24, 0x24, 0x40))
add_text_box(slide, Inches(1.2), Inches(6.1), Inches(10.8), Inches(0.5),
             "Overall automation potential: 60-70% reduction with proper tooling + agent workflows",
             font_size=16, color=GREEN_ACCENT, bold=True, align=PP_ALIGN.CENTER)

add_bottom_bar(slide)
add_slide_number(slide, 12, TOTAL_SLIDES)

# ============================================================
# SLIDE 13: Summary / Closing
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.04), ACCENT_BLUE)
add_shape(slide, Inches(0), Inches(0.04), SLIDE_W, Inches(0.02), ACCENT_TEAL)

add_text_box(slide, Inches(1.5), Inches(1.0), Inches(10), Inches(0.8),
             "VModelWorkflow", font_size=42, color=LIGHT_TEXT, bold=True,
             align=PP_ALIGN.CENTER, font_name="Calibri Light")

add_shape(slide, Inches(5.5), Inches(1.9), Inches(2.3), Inches(0.03), ACCENT_BLUE)

attributes = [
    ("Modular", "Adopt what you need, each component works independently", ACCENT_BLUE),
    ("Incremental", "Works on legacy code, module-by-module", ACCENT_TEAL),
    ("Domain-Agnostic", "Same core across aviation, automotive, medical, railway", ORANGE),
    ("AI-Native", "Designed for mid-senior engineers orchestrating AI agents", RGBColor(0xAA, 0x66, 0xCC)),
    ("Documentation-Proven", "Every claim backed by exhaustive best practices documentation", GREEN_ACCENT),
]

attr_start_y = Inches(2.4)
for i, (title, desc, color) in enumerate(attributes):
    y = attr_start_y + i * Inches(0.7)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.5), y + Inches(0.12),
                                  Inches(0.18), Inches(0.18))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    add_text_box(slide, Inches(2.9), y, Inches(2.5), Inches(0.4),
                 title, font_size=17, color=color, bold=True)
    add_text_box(slide, Inches(5.5), y + Inches(0.02), Inches(5.5), Inches(0.4),
                 desc, font_size=14, color=SUBTLE_TEXT)

# Closing quote
add_shape(slide, Inches(2), Inches(6.0), Inches(9.3), Inches(0.04), ACCENT_BLUE)
add_text_box(slide, Inches(1.5), Inches(6.3), Inches(10.3), Inches(0.7),
             "The V-model was designed for human discipline.\nIt turns out to be the perfect structure for AI discipline too.",
             font_size=18, color=ACCENT_TEAL, bold=True, align=PP_ALIGN.CENTER)

add_bottom_bar(slide)
add_slide_number(slide, 13, TOTAL_SLIDES)

# Save
output_path = "/home/stefanus/repos/VModelWorkflow/VModelWorkflow_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to {output_path}")
print(f"Total slides: {TOTAL_SLIDES}")
